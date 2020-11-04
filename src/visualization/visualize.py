import pandas as pd 
import geopandas as gp
import matplotlib.pyplot as plt 
import seaborn as sns
import math 
from ..utils import VIS_DIR

demographics=[
    'median_income' ,
    'pc_employed',
    'median_rent',
    'gini_index',
    'percent_income_spent_on_rent',
    'pc_occupied_homes' ,
    'pc_white',
    'pc_black',
    'pc_asian',
    'pc_hispanic'
]

map_crs= 'epsg:3857'
        
def map_self_initiated(city, ax=None, call_type=None, norm_by=None, year=None, vmin=None, vmax=None, scheme=None):
    if(not ax):
        fig = plt.figure()
        ax = fig.add_subplot(111)
    
    try:
        units = {
            'total': 'fraction',
            'area' : 'per m^{2}',
            'capita' : 'per capita'
        }

        title  = 'Self Initated {}'.format(units[norm_by])
        if(call_type):
            title = '{}: {}'.format(call_type, title)
        if(year):
            title = '{} - {}'.format(year,title)

        ax.set_title(title)
        city.self_initated_by_tract(norm_by=norm_by, call_type=call_type, year=None).plot(
                 column='Yes',
                 legend=True, 
                 vmin=vmin,
                 vmax=vmax, 
                 ax=ax,
                 scheme=scheme
        )
        ax.set_axis_off()
        ax.set_title('Officer initiated fraciton')
        return ax
    except:
        show_no_self_initiated_error(ax)
        return ax
    
def plot_cfs_breakdown(city,year=None,ax=None):
    ax= (city.filter_calls_by(year=year)
         .groupby('call_type')
         .count()['index']
         .sort_values(ascending=True)
         .plot(kind='barh',ax=ax, color=['#159BA3'] )
        )
    ax.set_title('{} - {} - CFS code'.format(city.BASE_NAME,year if year else 'All years' ))
    ax.set_xlabel('No of calls')
    ax.set_ylabel('')
    return ax        
        

def map_all_call_types(city, ax=None, norm_by=None, year=None, vmin=None, vmax=None, per_row=5):
    data = city.call_types_by_tract(norm_by=norm_by)
#     if(norm_by =='total'):
#         vmin = 0
#         vmax = 1
    no_cols = data.shape[0]
    columns = [col for col in data.columns if col not in ['geometry']]
    
    rows = math.ceil(len(columns)/per_row)
    fig = plt.figure(figsize=(5*per_row,5*rows))
    
    for index,column in enumerate(columns):
        ax = plt.subplot(rows,per_row,index+1)
        data.to_crs(map_crs).plot(column=column, vmin=vmin, vmax=vmax,ax=ax,legend=True, scheme='percentiles')
        ax.set_title(column,fontdict={'fontsize':'small'} )
        ax.set_axis_off()
        
    plt.suptitle('Calls per m^2')
    plt.tight_layout()

        
def plot_self_initated_by_call_type(city, ax=None, year=None):
    try:
        if(not ax):
            fig = plt.figure()
            ax = fig.add_subplot(111)
        data = city.self_initated_by_call_type(year=year)

        data.plot(kind='barh', stacked=True, label='Self Initiated',ax=ax, colors=['#159BA3','#F9BA16','#D1D4C9'])
        plt.xlabel('Fraction')
        plt.ylabel(None)
    #     box = ax.get_position()
    #     ax.set_position([box.x0, box.y0, box.width * 0.8, box.height])
    #     ax.legend(loc='center left', bbox_to_anchor=(1, 0.5))
        ax.legend(['Non self initiated', 'Self initiated','other'],loc='upper left')
        ax.set_xlabel('Fraction of calls')
        ax.set_ylabel('')
        ax.set_title('Officer initiated fraction by call type')
        ax.set_xlim(0,1)
        return ax
    except:
        show_no_self_initiated_error(ax)
        return ax 


def show_no_self_initiated_error(ax):
    ax.text(0.5,0.5,'No call type data', horizontalalignment='center', verticalalignment='center')
    ax.set_axis_off()
    return ax

def show_no_disposition_error(ax):
    ax.text(0.5,0.5,'No disposition data', horizontalalignment='center', verticalalignment='center')
    ax.set_axis_off()
    return ax

def plot_disposition_fraction_by_call_Type(city,ax=None, year=None):
    if(not ax):
        fig=plt.figure()
        ax= fig.add_subplot(111)
    try:
        data = city.disposition_fraction_by_call_type(year=year)
        (data.sort_values(by=city.ENFORCEMENT_VARIABLES)
            .plot(kind='barh', stacked=True, label='Disposition',ax=ax, color=['#1a15a3','#159BA3','#16120B','#f91693','#F9BA16','#E07323','#D1D4C9']))
    #     box = ax.get_position()
    #     ax.set_position([box.x0, box.y0, box.width * 0.8, box.height])
    #     ax.legend(loc='center left', bbox_to_anchor=(1, 0.5))
        ax.legend(loc='upper left')
        ax.set_xlabel("Fraction of calls")
        ax.set_title("Disposition fraction by CFS code")
        ax.set_ylabel('')
        return ax 
    except:
        show_no_disposition_error(ax)
        return ax
    
# def map_call_volume(city, ax=None, norm_by='None',year=None, call_type=None):
#     if(not ax):
#         fig = plt.figure()
#         ax = fig.add_subplot(111)
#     data = city.call_volumne_by_tract(year=year, call_type=call_type)
#     data.to_crs(map_crs).plot(column='calls', legend=True,ax=ax,scheme='quantiles')
#     ax.set_title('Calls per km^2')
#     ax.set_axis_off()
#     return ax
    
def map_call_type(city, ax=None, call_type=None, norm_by=None, year=None, vmin=None, vmax=None):
    
    if(not ax):
        fig = plt.figure()
        ax = fig.add_subplot(111)
        
    units = {
        'total': 'fraction',
        'area' : 'per m^{2}',
        'capita' : 'per capita'
    }

    city.call_type_by_tract(norm_by=norm_by, call_type=call_type).plot(
             column='No',
             legend=True, 
             vmin=vmin,
             vmax=vmax, 
             ax=ax 
    )
    ax.set_axis_off()
    return ax

def map_call_volume(city,norm_by=None, ax=None, year=None, call_type=None,scheme=None, vrange=None):
    if(not ax):
        fig = plt.figure()
        ax = fig.add_subplot(111)
    data = city.call_volume_by_tract(norm_by='capita',year=year, call_type=call_type)
    if(vrange):
        data.to_crs(map_crs).plot(column='calls', legend=True,ax=ax, scheme=scheme, vmin=vrange[0],vmax=vrange[1] )
    else:
        data.to_crs(map_crs).plot(column='calls', legend=True,ax=ax, scheme=scheme )
    ax.set_title('Calls per capita')
    ax.set_axis_off()
    return ax 

def plot_response_time_dist(city, call_type=None, year=None, ax=None):
    if(not ax):
        fig = plt.figure()
        ax = fig.add_subplot(111)
        
    data=city.filter_calls_by(
            call_type=call_type, 
            year=year
    ).response_time.dropna().div(60*60)
    
    sns.distplot(data, ax=ax, hist =True, color='#009aa6')
    ax.set_xlabel('Respose Time (h)')
    ax.set_ylabel('Call Fraction')
    ax.set_title('Distribution of response times (h)')
    return ax 

def map_median_response_time(city,call_type=None, ax=None, year=None,scheme=None):
    if(not ax):
        fig = plt.figure()
        ax = fig.add_subplot(111)
    data = city.median_response_time_by_tract(call_type=call_type, year=year)
    data.to_crs(map_crs).plot(column = 'response_time', ax=ax, legend=True, scheme=scheme)
    ax.set_title('Median response time (m)')
    ax.set_axis_off()
    return ax 

def plot_demographics_vs_call_volume(city,call_type,ax=None):
    pass
   
def plot_demographics_vs_disposition_fraction(city,call_type,ax=None):
    pass

def plot_demograpics_vs_self_initated(city,call_type,ax=None):
    pass
    
    
def plot_self_initated_by_disposition(city, call_type=None, year=None, ax=None):
    try:
        if(not ax):
            fig = plt.figure()
            ax = fig.add_subplot(111)
        data = city.self_initiated_by_disposition(call_type=call_type, year=year)
        data.sort_values(by='Yes').plot(kind='barh', ax=ax,stacked=True, color=['#FEBC16','#a6205c'], label='Self Initiated')
        ax.set_title('Officer initiated fraction of disposition')
        ax.legend(['Other', 'Officer initiated'], loc='upper left')
        ax.set_xlabel('Fraction of calls')
        ax.set_ylabel('')
    except:
        show_no_self_initiated_error(ax)
    
    return ax

def plot_disposition_counts(city,call_type=None, year=None,ax=None):
    if(not ax):
        fig = plt.figure()
        ax = fig.add_subplot(111)
    (city.disposition_counts(call_type=call_type, year=year)
        .plot(kind='barh',ax=ax, color='#009aa6')
    )
    ax.set_title('No of calls of disposition type')
    ax.set_xlabel('No of calls')
    ax.set_ylabel('')
    
    return ax


def map_enforcement_by_tract(city,call_type=None,year=None,ax=None,scheme=None):
    if(not ax):
        fig = plt.figure()
        ax = fig.add_subplot(111)
    try:
        (city.disposition_by_tract(call_type=call_type,year=year,norm_by='total')
             .assign( total_enforcement_action = lambda x: x[city.ENFORCEMENT_VARIABLES].sum(axis=1))
             .to_crs(map_crs)
             .plot(column='total_enforcement_action',legend=True, ax=ax, scheme=scheme)
        )
        ax.set_title('Enforcment Activity %')
        ax.set_axis_off()
        return ax
    except:
        show_no_disposition_error(ax)
        return ax

def make_demographics_report(city,cell_size=3, year=None, title=None, call_type=None, demos=None):
    if(not demos):
        demos=demographics 
        
    demo_labels  = {
       'percent_income_spent_on_rent': 'Percent income spent on rent',
       'pc_occupied_homes': 'Fraction of occupied homes',
       'pc_white' : 'Fraction of population white',
       'pc_black' : 'Fraction of population black',
       'pc_asian' : 'Fraction of population asian',
       'pc_hispanic' : 'Fraction of population hispanic',
       'median_income':'Median income $',
       'median_rent' : 'Median rent $',
       'pc_employed' : 'Fraction employed',
       'gini_index' : 'Gini Index'
    }
    
    grid_layout = (len(demos)+1, 4)
    fig = plt.figure(figsize=( (grid_layout[1]*cell_size*1.5, grid_layout[0]*cell_size)) )
    data = city.select_demographics(year=year,call_type=call_type)
      
    response_times = city.assign_demographics(city.median_response_time_by_tract(year=year, call_type=call_type))
    self_initatiated = city.assign_demographics(city.self_initated_by_tract(year=year,call_type=call_type))
    call_volumne  = city.assign_demographics(city.call_volume_by_tract(year=year,call_type=call_type))
    print('loaded data running demos')
    
    index= 0
    for demo in demos:
        print('running demo ', demo)
        
        index = index +1
        ax1 = plt.subplot(grid_layout[0], grid_layout[1], index)
        try:
            dispositions = city.assign_demographics(city.disposition_by_tract(year=year,call_type=call_type))
            dispositions = dispositions.assign( total_enforcement_action = lambda x: x[city.ENFORCEMENT_VARIABLES].sum(axis=1))
            sns.regplot(demo, 'total_enforcement_action', data=dispositions[dispositions[demo] > 0],ax=ax1, color ='#E47328')
            ax1.set_ylabel('Enforcement Activity')
            ax1.set_xlabel(demo_labels[demo])
        except:
            show_no_disposition_error(ax1)

        index = index +1
        ax2 = plt.subplot(grid_layout[0], grid_layout[1], index)
        sns.regplot(demo, 'calls' , data=call_volumne[call_volumne[demo]>0], ax=ax2, color='#2099A6')
        ax2.set_ylabel('Call Volume')
        ax2.set_xlabel(demo_labels[demo])

        index = index +1
        ax3 = plt.subplot(grid_layout[0], grid_layout[1], index)
        sns.regplot(demo, 'response_time' , data=response_times[response_times[demo]>0].assign(response_time = lambda x: x.response_time.div(60*60)), ax=ax3, color='#FEBC16')
        ax3.set_ylabel('Response Time (h)')
        ax3.set_xlabel(demo_labels[demo])
        
        try:
            index = index +1
            ax4 =plt.subplot(grid_layout[0], grid_layout[1], index)
            sns.regplot(demo, 'Yes' , data=self_initatiated[self_initatiated[demo]>0], ax=ax4, color ='#a6205c')
            ax4.set_ylabel('Officer initiated fraction')
            ax4.set_xlabel(demo_labels[demo])
        except:
            show_no_self_initiated_error(ax4)

    plt.suptitle("{} - {} - {}".format(city.BASE_NAME, year if year else 'all years' , call_type if call_type else 'all CFS codes'),
                 fontweight='bold', fontsize='x-large', color='#009aa6')
    plt.tight_layout(rect=[0, 0.03, 1, 0.94])
    plt.subplots_adjust(hspace=0.5)
    
def make_call_type_report(city, cell_size=8, year=None, title=None, call_type=None):
    grid_layout = (2,4)

    fig = plt.figure(figsize=( (grid_layout[1]*cell_size/1.5, grid_layout[0]*cell_size/1.5)) )
     #maps layout
    ax_calls_per_capita_map = plt.subplot2grid(grid_layout, (0,0),fig=fig)
    ax_self_initated_map = plt.subplot2grid(grid_layout, (0,1), fig=fig)
    ax_population_map = plt.subplot2grid(grid_layout, (1,0),fig=fig)
    ax_resonse_time_map = plt.subplot2grid(grid_layout, (1,1),fig=fig)
    
    #maps
    print('mapping call volune')
    map_call_volume(city,ax=ax_calls_per_capita_map, norm_by='capita', year=year,call_type=call_type)
    print("mapping self initated")
    map_self_initiated(city, norm_by='total',ax = ax_self_initated_map, year=year, call_type=call_type)
    print('mapping population')
    map_population(city,ax=ax_population_map)
    print("mapping response time ")
    map_median_response_time(city, ax=ax_resonse_time_map, year=year, call_type=call_type)
    
    ax_self_initiated_by_disposition = plt.subplot2grid(grid_layout, (0,2),colspan =1,fig=fig)
    ax_response_time_hist=plt.subplot2grid(grid_layout,(0,3))
    ax_enforcemnt_fraction =plt.subplot2grid(grid_layout,(1,2))
    ax_disposition_breakdown=plt.subplot2grid(grid_layout,(1,3))

    plot_self_initated_by_disposition(city,call_type=call_type, year=year, ax= ax_self_initiated_by_disposition)
    map_enforcement_by_tract(city,call_type=call_type, year=year, ax= ax_enforcemnt_fraction)
    plot_disposition_counts(city,call_type=call_type, year=year, ax= ax_disposition_breakdown)
    plot_response_time_dist(city,call_type=call_type, year=year, ax= ax_response_time_hist)
    
    plt.tight_layout(rect=[0, 0.03, 1, 0.94])
    if(title):
        plt.suptitle(title,fontweight='bold', fontsize='xx-large', color='#009aa6')
    else:
        plt.suptitle("{} - {} - {}".format(city.BASE_NAME,  year if year else 'All Years',call_type if call_type else 'All CFS codes'), fontweight='bold', fontsize='xx-large',color='#009aa6')

#     plt.subplots_adjust(hspace=0.5)


def make_simple_call_type_report(city, cell_size=8, year=None, title=None, call_type=None):
    plt.figure(figsize=(20,5))
    ax = plt.subplot(131)
    plot_self_initated_by_disposition(city,call_type=call_type, year=year, ax= ax)
    
    ax = plt.subplot(132)
    plot_disposition_counts(city,call_type=call_type, year=year, ax= ax)
    
    ax = plt.subplot(133)
    plot_response_time_dist(city,call_type=call_type, year=year, ax= ax)
    
    plt.tight_layout(rect=[0, 0.03, 1, 0.94])
    if(title):
        plt.suptitle(title,fontweight='bold', fontsize='xx-large', color='#009aa6')
    else:
        plt.suptitle("{} - {} - {}".format(city.BASE_NAME,  year if year else 'All Years',call_type if call_type else 'All CFS codes'), fontweight='bold', fontsize='xx-large',color='#009aa6')

#     plt.subplots_adjust(hspace=0.5)

def make_city_level_charts(city,cell_size=8, year=None,title=None, call_type=None):
    fig, axs = plt.subplots(2,2, figsize=(20,10))
    axs =axs.flatten()
    
    (city.filter_calls_by(year=year)
         .groupby('call_type')
         .count()['index']
         .sort_values(ascending=True)
         .plot(kind='barh',ax=axs[0], color=['#159BA3'] )
    )
    axs[0].set_title('CFS code')
    axs[0].set_xlabel('No of calls')
    axs[0].set_ylabel('')
    
    print('Plotting self initated by call type')
    plot_self_initated_by_call_type(city,ax=axs[1], year=year)
    
    print('plotting respone time distribution')
    plot_response_time_dist(city,ax=axs[2], year=year)
    plot_disposition_fraction_by_call_Type(city,ax=axs[3], year=year)
   
    if(title):
        plt.suptitle(title, fontweight='bold', fontsize='xx-large',color='#009aa6' )
    else:
        plt.suptitle("{} - {} - {}".format(city.BASE_NAME, call_type if call_type else 'All Calls', year if year else 'All Years'), fontweight='bold', fontsize='x-large',color='#009aa6')
        
    plt.tight_layout(rect=[0,0.03,1,0.94])
def make_city_level_maps(city,cell_size=8, year=None, title=None,call_type=None, scheme=None):
    
    if(scheme):
        legend_kwds={'loc': 'lower right'}
    
    fig, axs = plt.subplots(2, 3, figsize=(20,10))
    axs= axs.flatten()
    map_call_volume(city,ax=axs[0], 
                         norm_by='capita', 
                         year=year,
                         scheme=scheme, 
                         call_type=call_type)
    
    map_self_initiated(city, norm_by='total',
                             ax = axs[1], 
                             year=year, 
                             scheme=scheme, 
                             call_type=call_type)
    
    map_population(city,ax=axs[2], scheme=scheme)
    
    map_median_response_time(city, ax=axs[3], 
                                   year=year, 
                                   scheme=scheme, 
                                   call_type=call_type)
    
    map_enforcement_by_tract(city, ax=axs[4],
                                   scheme=scheme,
                                   year=year, 
                                   call_type=call_type)
    
    axs[5].set_axis_off()
    plt.tight_layout(rect=[0, 0.03, 1, 0.94])
    plt.subplots_adjust(hspace=0.0)

    pos = axs[3].get_position()
    axs[3].set_position([ pos.x0 + pos.width/2, pos.y0 , pos.width, pos.height])
    
    pos = axs[4].get_position()
    axs[4].set_position([ pos.x0 + pos.width/2, pos.y0  , pos.width, pos.height])
    
    
    if(title):
        plt.suptitle(title, fontweight='bold', fontsize='xx-large',color='#009aa6' )
    else:
        plt.suptitle("{} - {} - {}".format(city.BASE_NAME, call_type if call_type else 'All Calls', year if year else 'All Years'), fontweight='bold', fontsize='x-large',color='#009aa6')
        
    return axs
    
def make_city_level_report(city, cell_size=8, year=None, title=None, call_type=None):
    grid_layout = (2,4)

    fig = plt.figure(figsize=( (grid_layout[1]*cell_size/1.5, grid_layout[0]*cell_size/1.5)) )
    
    #maps layout
    ax_calls_per_capita_map = plt.subplot2grid(grid_layout, (0,0),fig=fig)
    ax_self_initated_map = plt.subplot2grid(grid_layout, (0,1), fig=fig)
    ax_population_map = plt.subplot2grid(grid_layout, (1,0),fig=fig)
    ax_resonse_time_map = plt.subplot2grid(grid_layout, (1,1),fig=fig)
    
    ax_call_breakdown = plt.subplot2grid(grid_layout, (0,2),colspan =1,fig=fig)
    ax_self_initated_by_call_type= plt.subplot2grid(grid_layout,(1,2), colspan =1, fig=fig)
    
    ax_response_time_hist=plt.subplot2grid(grid_layout,(0,3))
    ax_disposition_breakdown=plt.subplot2grid(grid_layout,(1,3))

    #maps
    print('mapping call volune ',year)
    map_call_volume(city,ax=ax_calls_per_capita_map, norm_by='capita', year=year)
    print("mapping self initated")
    map_self_initiated(city, norm_by='total',ax = ax_self_initated_map, year=year)
    print('mapping population')
    map_population(city,ax=ax_population_map)
    print("mapping response time ")
    map_median_response_time(city, ax=ax_resonse_time_map, year=year)
    
    print('Plotting call type breakdown')
    (city.filter_calls_by(year=year)
         .groupby('call_type')
         .count()['index']
         .sort_values(ascending=True)
         .plot(kind='barh',ax=ax_call_breakdown, color=['#159BA3'] )
    )
    ax_call_breakdown.set_title('CFS code')
    ax_call_breakdown.set_xlabel('No of calls')
    ax_call_breakdown.set_ylabel('')
    
    print('Plotting self initated by call type')
    plot_self_initated_by_call_type(city,ax=ax_self_initated_by_call_type, year=year)
    
    print('plotting respone time distribution')
    plot_response_time_dist(city,ax=ax_response_time_hist, year=year)
    plot_disposition_fraction_by_call_Type(city,ax=ax_disposition_breakdown, year=year)
   
    if(title):
        plt.suptitle(title, fontweight='bold', fontsize='xx-large',color='#009aa6' )
    else:
        plt.suptitle("{} - {} - {}".format(city.BASE_NAME, call_type if call_type else 'All Calls', year if year else 'All Years'), fontweight='bold', fontsize='x-large',color='#009aa6')
        
    plt.tight_layout(rect=[0, 0.03, 1, 0.94])
#     plt.subplots_adjust(hspace=0.5)
#     ax_self_initated_by_call_type.

def map_population(city,ax=None, scheme=None):
    if(not ax):
        fig =plt.figure()
        ax = fig.add_subplot(111)
    print('mapping population')
    (city.load_tracts().to_crs(map_crs)
         .plot(column="B01003_001E", legend=True,ax=ax,  scheme=scheme)
    )
    ax.set_title("Population")
    ax.set_axis_off()
    
    return ax 

def make_images_for_all(city, years=None):
    make_city_level_report(city)
    
    BASE_DIR = VIS_DIR /  'figures' / city.BASE_NAME 
    BASE_DIR.mkdir(parents=True, exist_ok=True)
    demographics_per_page = 3 
    demo_pages = range(0,math.ceil(len(demographics)/demographics_per_page))
    plt.savefig(str(BASE_DIR / 'all.png'))
    
    for i in demo_pages :
        demos = demographics[i*demographics_per_page:(i+1)*demographics_per_page]
        make_demographics_report(city,demos=demos)
        plt.savefig(str(BASE_DIR / 'all_demos_{}.png').format(i))
    if not years:
        years = city.processed_data.year.unique()
    print('Running for years ', years)
        
    for year in years:
        print('doing year', year ) 
        make_city_level_report(city,year=year)
        plt.savefig(str( BASE_DIR / 'all_{}.png'.format(year)))
        
        for i in demo_pages:
            demos = demographics[i*demographics_per_page:(i+1)*demographics_per_page]
            make_demographics_report(city,demos=demos,year=year)
            plt.savefig(str(BASE_DIR / 'all_{}_demos_{}.png').format(year,i))

    for crime in city.processed_data.call_type.unique():
        make_call_type_report(city,call_type =crime)
        if(crime == None):
            plt.savefig(str(BASE_DIR / 'other.png'))
        else:    
            plt.savefig(str(BASE_DIR / '{}.png'.format(crime.replace(" ","_").replace("/","_"))))
            
        for i in demo_pages:
            demos = demographics[i*demographics_per_page:(i+1)*demographics_per_page]
            make_demographics_report(city,call_type =crime,demos=demos)
            if(crime == None):
                plt.savefig(str(BASE_DIR / 'other_demos_{}.png'.format(i)))
            else:    
                plt.savefig(str(BASE_DIR / '{}_demos_{}.png'.format(crime.replace(" ","_").replace("/","_"),i)))

def make_spreadsheet_for_all(city, years=None):
    from pptx import Presentation
    from pptx.util import Inches
    

    left = Inches(0.44)
    top  = Inches(0.46)
    width = Inches(9)
    prs = Presentation()
    prs.slide_width=Inches(10)
    prs.slide_height= Inches(5.625)
    blank_slide_layout = prs.slide_layouts[6]
    
    demographics_per_page = 3 
    demo_pages = range(0,math.ceil(len(demographics)/demographics_per_page))

    BASE_DIR = VIS_DIR / 'figures' / city.BASE_NAME
    BASE_DIR.mkdir(parents=True, exist_ok=True)

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture(str(BASE_DIR / 'all.png'), left, top, width=width)
    
    for i in demo_pages:
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(str(BASE_DIR / 'all_demos_{}.png'.format(i)), left, top, width=width)
    
    if not years:
        years = city.processed_data.year.unique()
        
    for year in years:
        print('doing year', year ) 
        
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(str( BASE_DIR / 'all_{}.png'.format(year)), left, top,width=width)
        
        for i in demo_pages:
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(str( BASE_DIR / 'all_{}_demos_{}.png'.format(year,i)), left, top,width=width)
        
    for crime in city.processed_data.call_type.unique():
        print('doing cime, crime')
        try:
            slide = prs.slides.add_slide(blank_slide_layout)
            if(crime == None):
                slide.shapes.add_picture(str(BASE_DIR / 'other.png'), left, top, width=width)
            else:    
                slide.shapes.add_picture(str(BASE_DIR / '{}.png'.format(crime.replace(" ","_").replace("/","_"))), left, top, width=width)
            
            for i in demo_pages:
                slide = prs.slides.add_slide(blank_slide_layout)
                if(crime == None):
                    slide.shapes.add_picture(str(BASE_DIR / 'other_demos_{}.png').format(i), left, top, width=width)
                else:    
                    slide.shapes.add_picture(str(BASE_DIR / '{}_demos_{}.png'.format(crime.replace(" ","_").replace("/","_"),i)), left, top, width=width)
        except:
            print('issue with ', crime)
    prs.save(str(VIS_DIR / 'summaries_{}.pptx'.format(city.BASE_NAME)))

        
def plot_time_series(city,ax=None):
    if(not ax):
        fig = plt.figure()
        ax = fig.add_subplot(111)
    city.processed_data.groupby(city.processed_data.TimeCreate.dt.dayofyear).count()['index'].plot(ax=ax)
